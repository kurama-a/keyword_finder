import os
import pandas as pd
from docx import Document
from pptx import Presentation
import PyPDF2
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext

# Extraction du texte d'un fichier Word
def extract_text_from_docx(file_path, output_folder):
    doc = Document(file_path)
    text = ' '.join([para.text for para in doc.paragraphs])
    
    # Extraction des images
    for i, rel in enumerate(doc.part.rels.values()):
        if "image" in rel.target_ref:
            img_data = rel.target_part.blob
            img_filename = os.path.join(output_folder, f"{os.path.basename(file_path)}_image_{i}.png")
            with open(img_filename, "wb") as img_file:
                img_file.write(img_data)
            print(f"Image extraite: {img_filename}")
    
    return text

# Extraction du texte d'un fichier Excel
def extract_text_from_excel(file_path):
    data = pd.read_excel(file_path)
    return ' '.join(data.astype(str).values.flatten())

# Extraction du texte d'un fichier PowerPoint
def extract_text_from_pptx(file_path, output_folder):
    prs = Presentation(file_path)
    text = ''
    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                text += shape.text + ' '
            # Extraction des images
            if hasattr(shape, "image"):
                img = shape.image
                img_filename = os.path.join(output_folder, f"{os.path.basename(file_path)}_image_{i}_{j}.png")
                with open(img_filename, "wb") as img_file:
                    img_file.write(img.blob)
                print(f"Image extraite: {img_filename}")
    return text

# Extraction du texte d'un fichier PDF et des images
def extract_text_from_pdf(file_path, output_folder):
    text = ''
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text()

    # Convertir les pages PDF en images et traiter avec OCR
    images = convert_from_path(file_path)
    for i, image in enumerate(images):
        img_filename = os.path.join(output_folder, f"{os.path.basename(file_path)}_image_{i}.png")
        image.save(img_filename, 'PNG')
        print(f"Image extraite: {img_filename}")
    
    return text

# Extraction du texte d'un fichier TXT
def extract_text_from_txt(file_path):
    with open(file_path, 'r') as file:
        return file.read()

# Extraction du texte d'une image (OCR)
def extract_text_from_image(image_path):
    img = Image.open(image_path)
    text = pytesseract.image_to_string(img)
    return text

# Fonction pour détecter le type de fichier et extraire le texte correspondant
def extract_text(file_path, output_folder):
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == '.docx':
        return extract_text_from_docx(file_path, output_folder)
    elif ext == '.xlsx':
        return extract_text_from_excel(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path, output_folder)
    elif ext == '.pdf':
        return extract_text_from_pdf(file_path, output_folder)
    elif ext == '.txt':
        return extract_text_from_txt(file_path)
    elif ext in ['.jpg', '.jpeg', '.png', '.bmp']:
        return extract_text_from_image(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

# Fonction pour rechercher les mots-clés dans le texte extrait
def find_keywords_in_text(text, keywords):
    found_keywords = [keyword for keyword in keywords if keyword.lower() in text.lower()]
    return found_keywords

# Fonction pour rechercher les mots-clés dans plusieurs fichiers, y compris les sous-répertoires, et afficher les résultats
def search_keywords_in_files(directory_path, keywords, output_folder, output_file, result_text_widget):
    with open(output_file, 'w') as result_file:
        results = ""
        for root, dirs, files in os.walk(directory_path):
            for file_name in files:
                file_path = os.path.join(root, file_name)
                
                try:
                    # Extraire le texte du fichier
                    text = extract_text(file_path, output_folder)
                    
                    # Rechercher les mots-clés dans le texte
                    found_keywords = find_keywords_in_text(text, keywords)
                    
                    # Rechercher les mots-clés dans les images extraites
                    images_found_keywords = False
                    for img_file in os.listdir(output_folder):
                        if img_file.startswith(os.path.basename(file_path)):
                            img_path = os.path.join(output_folder, img_file)
                            if os.path.isfile(img_path):
                                ocr_text = extract_text_from_image(img_path)
                                found_keywords_in_image = find_keywords_in_text(ocr_text, keywords)
                                if found_keywords_in_image:
                                    images_found_keywords = True

                    # Si des mots-clés sont trouvés dans le fichier ou ses images, ajouter au résultat
                    if found_keywords or images_found_keywords:
                        result_str = f"Fichier : {file_path}\n"
                        if found_keywords:
                            result_str += f"Mots-cles trouves dans le texte : {', '.join(found_keywords)}\n"
                        if images_found_keywords:
                            result_str += f"Mots-cles trouves dans les images : {', '.join(found_keywords_in_image)}\n"
                        result_str += "-"*80 + "\n"
                        
                        result_file.write(result_str)
                        results += result_str
                    
                except Exception as e:
                    error_str = f"Erreur lors du traitement du fichier {file_name} : {e}\n"
                    result_file.write(error_str)
                    results += error_str
        
        if not results:
            results = "Aucun mot-clé trouvé dans les fichiers analysés."
            result_file.write(results)

        # Affichage dans le widget ScrolledText
        result_text_widget.config(state=tk.NORMAL)  # Permettre l'édition pour ajouter les résultats
        result_text_widget.delete(1.0, tk.END)  # Effacer l'ancien contenu
        result_text_widget.insert(tk.END, results)  # Ajouter les nouveaux résultats
        result_text_widget.config(state=tk.DISABLED)  # Désactiver l'édition après ajout

# Interface Graphique (GUI) avec Tkinter
def browse_directory():
    directory = filedialog.askdirectory()
    if directory:
        directory_path_entry.delete(0, tk.END)
        directory_path_entry.insert(0, directory)

def search_keywords():
    directory_path = directory_path_entry.get()
    keywords = keywords_entry.get().split(',')
    output_folder = 'images'
    output_file = 'results.txt'

    if not os.path.exists(directory_path):
        messagebox.showerror("Erreur", "Le répertoire spécifié n'existe pas.")
        return

    os.makedirs(output_folder, exist_ok=True)  # Créer le répertoire d'images s'il n'existe pas

    # Recherche des mots-clés et affichage des résultats
    search_keywords_in_files(directory_path, keywords, output_folder, output_file, result_text_widget)

# Création de la fenêtre principale
root = tk.Tk()
root.title("Recherche de mots-clés dans des fichiers")

# Champ pour le chemin du répertoire
tk.Label(root, text="Sélectionnez le répertoire :").grid(row=0, column=0, padx=10, pady=10)
directory_path_entry = tk.Entry(root, width=50)
directory_path_entry.grid(row=0, column=1, padx=10, pady=10)
tk.Button(root, text="Parcourir", command=browse_directory).grid(row=0, column=2, padx=10, pady=10)

# Champ pour les mots-clés
tk.Label(root, text="Mots-clés (séparés par des virgules) :").grid(row=1, column=0, padx=10, pady=10)
keywords_entry = tk.Entry(root, width=50)
keywords_entry.grid(row=1, column=1, padx=10, pady=10)

# Bouton de lancement de la recherche
tk.Button(root, text="Rechercher", command=search_keywords).grid(row=2, column=1, padx=10, pady=20)

# Zone de texte pour afficher les résultats avec une barre de défilement
result_text_widget = scrolledtext.ScrolledText(root, wrap=tk.WORD, width=80, height=20, font=("Arial", 10))
result_text_widget.grid(row=3, column=0, columnspan=3, padx=10, pady=10)
result_text_widget.config(state=tk.DISABLED)  # Désactiver l'édition pour la zone des résultats

# Lancer la boucle principale Tkinter
root.mainloop()

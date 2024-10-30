import os
import re
import pandas as pd
import argparse
from docx import Document
from pptx import Presentation
import PyPDF2
import win32com.client as win32
from concurrent.futures import ProcessPoolExecutor

PROGRESS_LOG_FILE = 'progress.log'
ERROR_LOG_FILE = 'error.log'

SUPPORTED_EXTENSIONS = {'.docx', '.doc', '.xlsx', '.pptx', '.pdf', '.txt', '.csv', '.odt', '.ods', '.odp', '.xls'}

# Fonction pour effacer les fichiers de log si nécessaire
def initialize_logs(overwrite_logs):
    if overwrite_logs:
        open(PROGRESS_LOG_FILE, 'w').close()
        open(ERROR_LOG_FILE, 'w').close()

# Fonction pour compter le nombre total de fichiers à scanner
def count_files(directory_path):
    total_files = sum(len([f for f in files if os.path.splitext(f)[-1].lower() in SUPPORTED_EXTENSIONS]) 
                      for _, _, files in os.walk(directory_path))
    return total_files

# Extraction du texte pour les différents formats de fichiers
def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = ' '.join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_doc(file_path):
    try:
        word = win32.gencache.EnsureDispatch("Word.Application")
        doc = word.Documents.Open(file_path)
        text = doc.Content.Text
        doc.Close()
        word.Quit()
        return text
    except ImportError:
        raise ImportError("pywin32 doit être installé pour extraire le texte d'un fichier .doc")

def extract_text_from_excel(file_path):
    data = pd.read_excel(file_path)
    return ' '.join(data.astype(str).values.flatten())

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + ' '
    return text

def find_urls(text):
    url_pattern = r'(https?://[^\s]+)'
    urls = re.findall(url_pattern, text)
    return urls

def extract_text_from_pdf(file_path):
    text = ''
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_text_from_csv(file_path):
    data = pd.read_csv(file_path, encoding='utf-8', sep=',')  # Modifier le séparateur si nécessaire
    text = ' '.join(data.astype(str).values.flatten())
    return text.lower() 

def extract_text_from_odt(file_path):
    import odf.opendocument
    from odf.text import P
    doc = odf.opendocument.load(file_path)
    text = ' '.join([str(p) for p in doc.getElementsByType(P)])
    return text

def extract_text_from_ods(file_path):
    data = pd.read_excel(file_path, engine='odf')
    return ' '.join(data.astype(str).values.flatten())

def extract_text_from_odp(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + ' '
    return text

def extract_text_from_xls(file_path):
    data = pd.read_excel(file_path)
    return ' '.join(data.astype(str).values.flatten())

# Fonction pour détecter le type de fichier et extraire le texte correspondant
def extract_text(file_path):
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.doc':
        return extract_text_from_doc(file_path)
    elif ext == '.xlsx':
        return extract_text_from_excel(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.txt':
        return extract_text_from_txt(file_path)
    elif ext == '.csv':
        return extract_text_from_csv(file_path)
    elif ext == '.odt':
        return extract_text_from_odt(file_path)
    elif ext == '.ods':
        return extract_text_from_ods(file_path)
    elif ext == '.odp':
        return extract_text_from_odp(file_path)
    elif ext == '.xls':
        return extract_text_from_xls(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

# Fonction pour rechercher les mots-clés et les URLs dans le texte extrait
def find_keywords_and_urls_in_text(text, keywords):
    text = text.lower()
    found_keywords = [keyword for keyword in keywords if keyword.lower() in text]
    urls = find_urls(text)
    return found_keywords, urls

# Fonction pour traiter un fichier et rechercher les mots-clés
def process_file(file_info, keywords):
    file_path, root_dir = file_info
    try:
        text = extract_text(file_path)
        found_keywords, urls = find_keywords_and_urls_in_text(text, keywords)
        
        if found_keywords:
            return {
                'Directory': os.path.abspath(root_dir), 
                'Filename': os.path.basename(file_path),
                'URLs': ';'.join(set(urls))
            }
    except Exception as e:
        ext = os.path.splitext(file_path)[-1].lower()
        if ext in SUPPORTED_EXTENSIONS:
            with open(ERROR_LOG_FILE, 'a') as error_file:
                error_file.write(f"Error processing {file_path}: {e}\n")
    return None

# Fonction pour lire l'état de progression
def load_progress():
    if os.path.exists(PROGRESS_LOG_FILE):
        with open(PROGRESS_LOG_FILE, 'r') as f:
            return set(line.strip() for line in f)
    return set()

# Fonction pour enregistrer l'état de progression en mode 'append'
def save_progress(processed_files):
    with open(PROGRESS_LOG_FILE, 'a') as f:
        for file in processed_files:
            f.write(f"{file}\n")

# Fonction principale pour rechercher les mots-clés en parallèle avec suivi de progression
def search_keywords_in_files(directory_path, keywords, csv_output_file, batch_size=100):
    results_data = []
    processed_files = load_progress()
    file_paths = [
        (os.path.join(root_dir, file_name), root_dir) 
        for root_dir, _, files in os.walk(directory_path) 
        for file_name in files 
        if os.path.splitext(file_name)[-1].lower() in SUPPORTED_EXTENSIONS
        and os.path.join(root_dir, file_name) not in processed_files
    ]

    total_files = len(file_paths)
    print(f"Nombre total de fichiers à scanner : {total_files}")

    current_file = 0
    with ProcessPoolExecutor() as executor:
        futures = []
        batch = []
        
        # Traiter les fichiers en lot pour sauvegarder périodiquement
        for file_info in file_paths:
            batch.append(file_info)
            current_file += 1
            print(f"Progression : {current_file}/{total_files}", end='\r')
            
            if len(batch) >= batch_size:
                futures.append(executor.submit(process_batch, batch, keywords))
                batch = []

        # Traiter les fichiers restants
        if batch:
            futures.append(executor.submit(process_batch, batch, keywords))
        
        # Collecte des résultats
        for future in futures:
            batch_results, batch_files = future.result()
            results_data.extend(batch_results)
            processed_files.update(batch_files)
            save_progress(batch_files)

    # Sauvegarder les résultats finaux dans le fichier CSV en mode 'append'
    if results_data:
        df = pd.DataFrame(results_data)
        df.to_csv(csv_output_file, index=False, mode='a', header=not os.path.exists(csv_output_file))

    print(f"\nScan terminé. Résultats sauvegardés dans {csv_output_file}")

# Fonction pour traiter un lot de fichiers
def process_batch(file_batch, keywords):
    batch_results = []
    batch_files = []

    for file_info in file_batch:
        result = process_file(file_info, keywords)
        if result:
            batch_results.append(result)
            batch_files.append(file_info[0])

    return batch_results, batch_files

# Fonction pour charger les mots-clés à partir d'un fichier
def load_keywords(keyword_file):
    with open(keyword_file, 'r') as file:
        return file.read().strip().split(',')

# Fonction principale
def main():
    parser = argparse.ArgumentParser(description="Recherche de mots-clés dans des fichiers.")
    parser.add_argument('directory', type=str, help="Chemin du répertoire à scanner.")
    parser.add_argument('keyword_file', type=str, help="Fichier contenant les mots-clés séparés par des virgules.")
    parser.add_argument('--csv_output_file', type=str, default='keyword_search_results.csv', help="Fichier CSV pour enregistrer les résultats.")
    parser.add_argument('--overwrite_logs', action='store_true', help="Réécrire les logs existants.")
    args = parser.parse_args()

    # Initialiser les logs selon l'option overwrite_logs
    initialize_logs(args.overwrite_logs)

    # Charger les mots-clés depuis le fichier
    keywords = load_keywords(args.keyword_file)

    # Recherche des mots-clés en parallèle et enregistrement des résultats
    search_keywords_in_files(args.directory, keywords, args.csv_output_file)

if __name__ == "__main__":
    main()
